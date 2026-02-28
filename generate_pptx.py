from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "市場開發價值提案 — Sales Copilot"
    subtitle.text = "和亞智慧（7825）OT 事業部\n提案人：沈文煜\n2026-02-25"

    # Slide 2: Problem & Challenge
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "痛點與挑戰：設備銷售的死穴"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "情報滯後導致陷入「規格戰」"
    
    p = tf.add_paragraph()
    p.text = "• 錯過 NPI 到 MP 的黃金開窗期 (4-8 週)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 等到客戶開出正式 Spec，只能淪為 Vendor 2 比價"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 傳統陌生開發效率低，無法精準打擊技術痛點"
    p.level = 1

    # Slide 3: Solution 1 - Market Radar
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "解法一：Market Radar (市場雷達)"
    tf = slide.placeholders[1].text_frame
    tf.text = "自動化情報追蹤，極限壓縮 Time-to-Insight"
    
    p = tf.add_paragraph()
    p.text = "• 自動爬取 CapEx 公告、環評擴廠文件、LinkedIn 職缺異動"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 監控目標：玉晶光、揚明光、禾賽、穩懋等關鍵供應鏈"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 價值：搶在開標前 2 個月，掌握客戶擴張動向"
    p.level = 1

    # Slide 4: Solution 2 - Pain Point Extractor
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "解法二：Pain Point Extractor (技術痛點提取)"
    tf = slide.placeholders[1].text_frame
    tf.text = "從「賣規格」轉向「賣解法」"
    
    p = tf.add_paragraph()
    p.text = "• AI 分析專利、法說會逐字稿、技術週報"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 鎖定痛點：Pancake 鏡片良率、VCSEL 點雲雜訊、AOI 誤判率"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 價值：開發信直指 RD 技術卡點，創造顧問式銷售信任感"
    p.level = 1

    # Slide 5: Solution 3 - BANT Scorer (Signal-driven)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "解法三：BANT Scorer (商機成熟度預評分)"
    tf = slide.placeholders[1].text_frame
    tf.text = "基於「代理數據」先行量化，優化開發優先序"
    
    p = tf.add_paragraph()
    p.text = "• Budget: 由 CapEx / 擴產公告 / 設備預算推估"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Need: 由 NPI 訊號 / 專利佈局 / 技術痛點匹配度推估"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Timeline: 由 專案週期 / 量產時程窗口推估"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 價值：真正互動前的「預篩選」，確保主管頻寬投在 A 級目標"
    p.level = 1

    # Slide 6: Commitment
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "即戰力承諾：Day 1 帶槍投靠"
    tf = slide.placeholders[1].text_frame
    tf.text = "入職後的開發藍圖"
    
    p = tf.add_paragraph()
    p.text = "• Week 1: 匯入 OT 事業部現有名單，啟動 Market Radar"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Week 2: 產出 LiDAR 與 AR/VR 關鍵客戶技術痛點分析"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Week 4: 建立「商機優先序看板 (Signal-driven Pipeline)」"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "目標：透過情報驅動，建立具備 BANT 雛形的實戰漏斗"
    p.level = 0

    output_path = "2.25市場開發價值提案_SalesCopilot_GEMINI.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()

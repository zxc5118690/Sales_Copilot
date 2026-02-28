from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import sys

def create_presentation():
    prs = Presentation()

    # Define common paths
    base_path = "/Users/shenwenyu/.gemini/antigravity/brain/32ba34cc-f257-41be-b56e-7e7f98ac4421"
    screenshots = {
        "Dashboard": os.path.join(base_path, "screenshot_dashboard_1771725612358.png"),
        "Market Radar": os.path.join(base_path, "screenshot_signals_1771725618456.png"),
        "Pain Extractor": os.path.join(base_path, "screenshot_pains_1771725624890.png"),
        "Outreach": os.path.join(base_path, "screenshot_outreach_1771725631143.png"),
        "Pipeline": os.path.join(base_path, "screenshot_pipeline_1771725637760.png")
    }

    def add_title_slide(title, subtitle):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title
        subtitle_shape.text = subtitle

    def add_content_slide(title, bullet_points, image_path=None):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        # Add bullets
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
            p.level = 0
            p.font.size = Pt(18)

        if image_path and os.path.exists(image_path):
            # Adjust body width to 4 inches (half slide)
            body_shape.width = Inches(4)
            # Add picture on the right side
            left = Inches(4.5)
            top = Inches(1.5)
            # Set height or width, keep aspect ratio
            # Let's target a height that fits well
            slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    # 1. 封面
    add_title_slide("和亞智慧 OT 事業部\nAI 賦能市場開發提速方案", "提案人：沈文煜\nSales Copilot × 智慧製造開路先鋒")

    # 2. 部門挑戰
    add_content_slide("OT 部門現狀與挑戰", [
        "Q3 案場開發壓力 vs Q4 營收實現目標",
        "AR/VR 與 LiDAR 高技術門檻市場開發耗時",
        "主管心力分配：技術決策 vs 陌生開發過濾",
        "Time-to-Insight：從市場情報到行動轉化的延遲"
    ])

    # 3. 解決方案：Sales Copilot
    add_content_slide("解決方案：Sales Copilot 核心價值", [
        "Market Radar：自動捕捉客戶 CapEx 與 NPI 訊號",
        "Pain Extractor：針對 RD 痛點自動生成開發切入點",
        "BANT Scorer：量化商商機成熟度，優化主管介入時機",
        "Pipeline Dashboard：實時掌握部門業務進度"
    ], screenshots["Dashboard"])

    # 4. 市場雷達 (Signals)
    add_content_slide("模組 1：市場雷達 (Market Radar)", [
        "即時追蹤玉晶光、揚明光、同欣電等擴產公告",
        "自動識別 104 招聘與法說會逐字稿資訊",
        "主動推送採購視窗開啟訊號 (Signal Score)",
        "減少 80% 人工搜尋與整理時間"
    ], screenshots["Market Radar"])

    # 5. 痛點提取 (Pains)
    add_content_slide("模組 2：痛點提取 (Pain Extractor)", [
        "分析 RD / NPI 在製程中的具體技術瓶頸",
        "自動生成 '對位' 方案的專業話術",
        "從專利公開資料中萃取潛在技術卡點",
        "大幅提升 Cold Email 的開信率與回覆精度"
    ], screenshots["Pain Extractor"])

    # 6. 訊息草稿 (Outreach)
    add_content_slide("模組 3：高轉化訊息生成 (Outreach)", [
        "不再是罐頭信：基於 Pain Profile 的個性化草稿",
        "自動生成 Email 與 LinkedIn 同步開發腳本",
        "CTA (Call to Action) 定向引導至技術 Demo",
        "確保每一封信都具備顧問級的專業度"
    ], screenshots["Outreach"])

    # 7. Pipeline 與 BANT
    add_content_slide("數據化管理：Pipeline & BANT", [
        "自動化 BANT 評分：Budget, Authority, Need, Timeline",
        "商機分級系統 (A/B/C)：精確引導主管時間投入",
        "入職 Day 1 即可運行的自動化專案看板",
        "決策透明度：主管隨時掌握每一案場的目前深度"
    ], screenshots["Pipeline"])

    # 8. 執行藍圖與承諾
    add_content_slide("入職 90 天計畫與承諾", [
        "Day 1-7：匯入 50+ 目標客戶，輸出首份雷達報告",
        "Day 14：完成 Top 10 客戶痛點分析與初版寄送",
        "Day 30：建立完整 BANT 分級 Pipeline",
        "Day 90：篩選出 3-5 組高品質專案供主管介入成交"
    ])

    save_path = "/Users/shenwenyu/Desktop/wencode/Interview-newsmart-salessoltion/市場開發價值提案_SalesCopilot.pptx"
    prs.save(save_path)
    print(f"Presentation saved to: {save_path}")

if __name__ == "__main__":
    create_presentation()
